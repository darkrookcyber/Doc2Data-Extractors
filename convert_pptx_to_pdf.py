"""
Tool: Convert .pptx file to .pdf
Requirements: Must install pythin-pptx and reportlab libraries -
1. Install the required packages by running:

```bash
pip install python-pptx reportlab
```

To Use:
To run the script from the terminal or command prompt:

- For a single `.pptx` to `.pdf` conversion: `python3 pptx_to_pdf.py [input filename].pptx 
[output filename].pdf`
- For no input file specified, and the script will ask for an input file during 
runtime: `python pptx_to_pdf.py`




import os
from pptx import Presentation
import reportlab.lib.pagesizes as pagesizes
from reportlab.platypus import Image as PlatypusImage, Paragraph
from reportlab.pdfbase import pdfbase, PdfFileReader
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib import colors
from reportlab.lib.units import mm

def convert_pptx_to_pdf(input_file, output_file):
    prs = Presentation(input_file)
    
    # Create a PDF document with ReportLab's PdfFileWriter
    pdf_writer = pdfbase.PdfBase()
    pdf_writer.builder.open()

    # Set up the page size and font
    page = Page(pagesizes.letter)
    page.buildTrimBox = mm(10.5) * 72
    page.mediaBox = pagesizes.letter
    c = pdf_writer.newPage(page)
    TTFont.RegisterFont(name='Arial', fnfile='arial.ttf')
    c.setFont('Arial', 12)

    # Convert slides to images and add them to the PDF
    for slide in prs.slides:
        img = PlatypusImage(slide.images[0])
        p = Paragraph(' ', width=img.width, height=img.height)
        p.image = img
        c.drawOn(p, 0, 0)
        image_file = PdfFileReader(slide.images[0])
        page = pdf_writer.newPage()
        page.mergePage(page)
        pdf_writer.addImage(image_file.getImageData(), pagesizes.letter[0], 
pagesizes.letter[1])

    # Save and close the PDF document
    pdf_writer.builder.saveFiles([output_file])
    pdf_writer.close()

if __name__ == '__main__':
    input_file = None
    
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
    
    if not input_file:
        print("Please provide a PPTX file as an argument.")
        print("Usage: python pptx_to_pdf.py [input_pptx_file.pptx] [output_pdf_file.pdf]")
        
    else:
        output_file = "output.pdf" if len(sys.argv) < 3 else sys.argv[2]
        convert_pptx_to_pdf(input_file, output_file)
        
"""
        
import argparse
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

def pptx_to_pdf(input_file_path, output_file_path):
    """
    Converts a .pptx file to .pdf format with improved text handling to avoid overlap,
    and allows specifying input and output files via command line or runtime prompt.

    Args:
    input_file_path (str): The path to the .pptx input file.
    output_file_path (str): The path to save the output .pdf file.
    """
    presentation = Presentation(input_file_path)
    c = canvas.Canvas(output_file_path, pagesize=letter)
    width, height = letter

    for slide in presentation.slides:
        y_position = height - 50
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    c.drawString(50, y_position, run.text)
                    y_position -= 20
                y_position -= 10
            if y_position < 50:
                c.showPage()
                y_position = height - 50
        c.showPage()
    c.save()

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert a .pptx document to .pdf format.")
    parser.add_argument("-i", "--input", help="The path to the .pptx document to be converted.")
    parser.add_argument("-o", "--output", default="output.pdf", help="The path to save the output .pdf file. Defaults to 'output.pdf'.")
    args = parser.parse_args()

    input_file_path = args.input
    if not input_file_path:
        input_file_path = input("Please enter the path to the .pptx document: ")

    output_file_path = args.output
    pptx_to_pdf(input_file_path, output_file_path)